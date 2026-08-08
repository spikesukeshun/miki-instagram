"""
make_report_pptx.py
────────────────────────────────────────────────────
insight_report_20260414.md の分析データを
視覚的に見やすいPowerPointに変換する。

実行:
  python3 make_report_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import json, os, requests
from datetime import datetime, timezone, timedelta

from load_env import load_from_zshrc
load_from_zshrc()

TOKEN      = os.getenv("INSTAGRAM_ACCESS_TOKEN")
ACCOUNT_ID = os.getenv("INSTAGRAM_BUSINESS_ACCOUNT_ID")
API_BASE   = "https://graph.facebook.com/v19.0"

# ── カラーパレット（エレガント・エステ調） ──────────────
C_BG_DARK  = RGBColor(0x1C, 0x1C, 0x1C)   # ほぼ黒（タイトルBG）
C_BG_LIGHT = RGBColor(0xFA, 0xF8, 0xF5)   # オフホワイト（本文BG）
C_GOLD     = RGBColor(0xC9, 0xA9, 0x6E)   # シャンパンゴールド
C_GOLD_DARK= RGBColor(0x8B, 0x69, 0x1A)   # ダークゴールド
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK    = RGBColor(0x1C, 0x1C, 0x1C)
C_GRAY     = RGBColor(0x6B, 0x6B, 0x6B)
C_LGRAY    = RGBColor(0xE8, 0xE4, 0xDE)   # ライトグレー（罫線）
C_GREEN    = RGBColor(0x4A, 0x8C, 0x5C)   # アクセント緑（ポジティブ）
C_RED      = RGBColor(0xB5, 0x4A, 0x4A)   # アクセント赤（警告）
C_BLUE     = RGBColor(0x4A, 0x6E, 0x9A)   # アクセント青

# スライドサイズ: ワイド (13.33 x 7.5 inch)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

JST = timezone(timedelta(hours=9))

OUTPUT = f"instagram_insight_{datetime.now().strftime('%Y%m%d')}.pptx"


# ─────────────────────────────────────────────────────
# ユーティリティ
# ─────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=Pt(14), bold=False, color=C_BLACK,
                 align=PP_ALIGN.LEFT, wrap=True, font_name="Hiragino Kaku Gothic Pro"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_label_value(slide, label, value, x, y, w,
                    label_color=C_GRAY, value_color=C_BLACK,
                    label_size=Pt(11), value_size=Pt(28), unit=""):
    """ラベル + 大きな数値を縦に並べる"""
    add_text_box(slide, label, x, y, w, Inches(0.35),
                 font_size=label_size, color=label_color, align=PP_ALIGN.CENTER)
    add_text_box(slide, str(value) + unit, x, y + Inches(0.3), w, Inches(0.7),
                 font_size=value_size, bold=True, color=value_color, align=PP_ALIGN.CENTER)


def add_bar(slide, x, y, bar_w, bar_h, value, max_value,
            fill_color=C_GOLD, bg_color=C_LGRAY):
    """水平バーグラフ"""
    add_rect(slide, x, y, bar_w, bar_h, fill_color=bg_color)
    ratio = min(value / max_value, 1.0) if max_value > 0 else 0
    if ratio > 0:
        add_rect(slide, x, y, int(bar_w * ratio), bar_h, fill_color=fill_color)


def add_divider(slide, x, y, w, color=C_GOLD, thickness=Pt(1.5)):
    line = slide.shapes.add_shape(1, x, y, w, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def slide_bg(slide, color=C_BG_LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ─────────────────────────────────────────────────────
# スライド1: タイトル
# ─────────────────────────────────────────────────────
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide_bg(slide, C_BG_DARK)

    # 金のアクセントライン（上）
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=C_GOLD)

    # メインタイトル
    add_text_box(slide, "Instagram インサイト分析レポート",
                 Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
                 font_size=Pt(40), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # サブタイトル
    add_text_box(slide, "2026年4月14日  |  全391投稿分析（ライブAPIデータ）",
                 Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.6),
                 font_size=Pt(18), color=C_GOLD, align=PP_ALIGN.CENTER)

    # 区切り線
    add_divider(slide, Inches(3.5), Inches(3.8), Inches(6.3), color=C_GOLD)

    # アカウント情報
    add_text_box(slide, "MIKIのエステInstagramアカウント",
                 Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.5),
                 font_size=Pt(16), color=C_GRAY, align=PP_ALIGN.CENTER)

    # 金のアクセントライン（下）
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_color=C_GOLD)


# ─────────────────────────────────────────────────────
# スライド2: アカウントサマリー（KPI）
# ─────────────────────────────────────────────────────
def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, C_BG_LIGHT)

    # ヘッダー
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=C_BG_DARK)
    add_text_box(slide, "アカウント概要", Inches(0.5), Inches(0.25), Inches(8), Inches(0.6),
                 font_size=Pt(24), bold=True, color=C_WHITE)
    add_text_box(slide, "2026年4月14日時点", Inches(10.5), Inches(0.35), Inches(2.5), Inches(0.4),
                 font_size=Pt(12), color=C_GOLD, align=PP_ALIGN.RIGHT)

    # KPIカード × 4
    cards = [
        ("総投稿数", "391", "件", C_GOLD_DARK),
        ("最高いいね数", "451", "件", C_GREEN),
        ("最高リーチ数", "937", "人", C_BLUE),
        ("最高保存数", "14", "件", C_GOLD),
    ]
    card_w = Inches(2.8)
    card_gap = Inches(0.3)
    card_x_start = Inches(0.55)

    for i, (label, value, unit, color) in enumerate(cards):
        cx = card_x_start + i * (card_w + card_gap)
        cy = Inches(1.5)
        ch = Inches(2.0)
        # カードBG
        add_rect(slide, cx, cy, card_w, ch, fill_color=C_WHITE,
                 line_color=C_LGRAY, line_width=Pt(1))
        # カラーバー（上部）
        add_rect(slide, cx, cy, card_w, Inches(0.1), fill_color=color)
        # ラベル
        add_text_box(slide, label, cx + Inches(0.1), cy + Inches(0.25), card_w - Inches(0.2), Inches(0.4),
                     font_size=Pt(12), color=C_GRAY, align=PP_ALIGN.CENTER)
        # 数値
        add_text_box(slide, value, cx + Inches(0.1), cy + Inches(0.6), card_w - Inches(0.2), Inches(0.85),
                     font_size=Pt(44), bold=True, color=color, align=PP_ALIGN.CENTER)
        # 単位
        add_text_box(slide, unit, cx + Inches(0.1), cy + Inches(1.45), card_w - Inches(0.2), Inches(0.35),
                     font_size=Pt(14), color=C_GRAY, align=PP_ALIGN.CENTER)

    # 補足テキスト
    notes = [
        ("📌 いいね最高投稿（2023-02-01）", "自己語り投稿「今日は誰も興味のない、自分語りです（笑）」"),
        ("📌 リーチ最高投稿（2024-06-18）", "プライベート投稿「久々のプライベート😊」"),
        ("📌 保存最高投稿（2023-02-01）", "同上の自己語り投稿（いいね・リーチ・保存すべてトップ）"),
    ]
    for j, (title, body) in enumerate(notes):
        ny = Inches(3.8) + j * Inches(0.9)
        add_rect(slide, Inches(0.5), ny, Inches(12.3), Inches(0.75),
                 fill_color=C_WHITE, line_color=C_LGRAY, line_width=Pt(0.5))
        add_text_box(slide, title, Inches(0.7), ny + Inches(0.05), Inches(3.5), Inches(0.35),
                     font_size=Pt(12), bold=True, color=C_GOLD_DARK)
        add_text_box(slide, body, Inches(4.2), ny + Inches(0.05), Inches(8.4), Inches(0.35),
                     font_size=Pt(12), color=C_BLACK)


# ─────────────────────────────────────────────────────
# スライド3: テーマ別パフォーマンス
# ─────────────────────────────────────────────────────
def slide_theme_perf(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, C_BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=C_BG_DARK)
    add_text_box(slide, "テーマ別パフォーマンス", Inches(0.5), Inches(0.25), Inches(8), Inches(0.6),
                 font_size=Pt(24), bold=True, color=C_WHITE)
    add_text_box(slide, "平均いいね数・投稿比率", Inches(10), Inches(0.35), Inches(3), Inches(0.4),
                 font_size=Pt(12), color=C_GOLD, align=PP_ALIGN.RIGHT)

    # データ
    themes = [
        ("ライフスタイル\n・自己語り", 29.9, 74, 19.1, C_GREEN),
        ("ブライダル",                20.3, 30,  7.7, C_BLUE),
        ("ご褒美エステ",              17.0,  7,  1.8, C_GOLD),
        ("メニュー・\nサービス",      11.8, 211, 54.0, C_RED),
        ("未分類",                    18.2, 69, 17.6, C_GRAY),
    ]

    bar_area_x = Inches(0.5)
    bar_area_y = Inches(1.4)
    bar_h = Inches(0.7)
    bar_gap = Inches(0.45)
    label_w = Inches(2.0)
    bar_max_w = Inches(6.0)
    max_likes = 30.0

    for i, (name, avg_likes, count, pct, color) in enumerate(themes):
        row_y = bar_area_y + i * (bar_h + bar_gap)

        # テーマ名
        add_text_box(slide, name, bar_area_x, row_y, label_w, bar_h,
                     font_size=Pt(13), bold=True, color=C_BLACK, align=PP_ALIGN.RIGHT, wrap=True)

        # バーグラフ
        bx = bar_area_x + label_w + Inches(0.15)
        add_bar(slide, bx, row_y + Inches(0.15), bar_max_w, bar_h - Inches(0.3),
                avg_likes, max_likes, fill_color=color)

        # 数値ラベル
        bar_end = bx + bar_max_w
        add_text_box(slide, f"平均 {avg_likes}いいね",
                     bar_end + Inches(0.1), row_y + Inches(0.1), Inches(1.5), bar_h - Inches(0.2),
                     font_size=Pt(13), bold=True, color=color)

        # 投稿数 + 割合
        add_text_box(slide, f"{count}投稿（{pct}%）",
                     bar_end + Inches(1.65), row_y + Inches(0.1), Inches(1.5), bar_h - Inches(0.2),
                     font_size=Pt(12), color=C_GRAY)

    # 凡例 / インサイト
    insight_y = Inches(5.65)
    add_divider(slide, Inches(0.5), insight_y, Inches(12.3), color=C_GOLD)
    add_text_box(slide, "💡  KEY INSIGHT",
                 Inches(0.5), insight_y + Inches(0.1), Inches(2.5), Inches(0.4),
                 font_size=Pt(13), bold=True, color=C_GOLD_DARK)
    add_text_box(slide,
                 "自己語り投稿のいいね平均はメニュー告知の 2.5倍。"
                 "しかしメニュー告知が投稿全体の 54% を占める — 比率の見直しが急務。",
                 Inches(3.0), insight_y + Inches(0.1), Inches(10.0), Inches(0.5),
                 font_size=Pt(13), color=C_BLACK)


# ─────────────────────────────────────────────────────
# スライド4: コンテンツタイプ別
# ─────────────────────────────────────────────────────
def slide_type_perf(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, C_BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=C_BG_DARK)
    add_text_box(slide, "コンテンツタイプ別パフォーマンス", Inches(0.5), Inches(0.25), Inches(9), Inches(0.6),
                 font_size=Pt(24), bold=True, color=C_WHITE)

    types = [
        ("🎬  リール動画", 67, 17.1, 18.6, C_GREEN),
        ("🖼  カルーセル", 260, 66.5, 17.8, C_GOLD),
        ("📷  単体画像",   64, 16.4, 12.6, C_GRAY),
    ]

    # カード型レイアウト
    card_w = Inches(3.8)
    gap    = Inches(0.3)
    start_x = Inches(0.5)
    card_y  = Inches(1.4)
    card_h  = Inches(3.8)

    for i, (label, count, pct, avg, color) in enumerate(types):
        cx = start_x + i * (card_w + gap)
        add_rect(slide, cx, card_y, card_w, card_h, fill_color=C_WHITE,
                 line_color=C_LGRAY, line_width=Pt(1))
        # 上部カラーバー
        add_rect(slide, cx, card_y, card_w, Inches(0.12), fill_color=color)

        # タイプ名
        add_text_box(slide, label, cx + Inches(0.2), card_y + Inches(0.25), card_w - Inches(0.4), Inches(0.5),
                     font_size=Pt(16), bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)

        add_divider(slide, cx + Inches(0.3), card_y + Inches(0.8), card_w - Inches(0.6), color=C_LGRAY, thickness=Pt(0.5))

        # 投稿数
        add_text_box(slide, "投稿数", cx + Inches(0.2), card_y + Inches(0.95), card_w - Inches(0.4), Inches(0.3),
                     font_size=Pt(11), color=C_GRAY, align=PP_ALIGN.CENTER)
        add_text_box(slide, f"{count}件", cx + Inches(0.2), card_y + Inches(1.2), card_w - Inches(0.4), Inches(0.5),
                     font_size=Pt(30), bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_box(slide, f"全投稿の {pct}%", cx + Inches(0.2), card_y + Inches(1.65), card_w - Inches(0.4), Inches(0.3),
                     font_size=Pt(12), color=C_GRAY, align=PP_ALIGN.CENTER)

        add_divider(slide, cx + Inches(0.3), card_y + Inches(2.05), card_w - Inches(0.6), color=C_LGRAY, thickness=Pt(0.5))

        # 平均いいね
        add_text_box(slide, "平均いいね数", cx + Inches(0.2), card_y + Inches(2.2), card_w - Inches(0.4), Inches(0.3),
                     font_size=Pt(11), color=C_GRAY, align=PP_ALIGN.CENTER)
        add_text_box(slide, f"{avg}件", cx + Inches(0.2), card_y + Inches(2.45), card_w - Inches(0.4), Inches(0.6),
                     font_size=Pt(36), bold=True, color=color, align=PP_ALIGN.CENTER)

    # リール補足
    add_rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.75),
             fill_color=RGBColor(0xE8, 0xF5, 0xEE), line_color=C_GREEN, line_width=Pt(1))
    add_text_box(slide,
                 "📌  リールは現在67本（17%）と少ないが、平均いいねはカルーセルと同等以上。"
                 "月1〜2本ペースで増やすとリーチ拡大が期待できる。"
                 "ストーリーズの過去データはAPI上取得不可（24時間で消滅）。",
                 Inches(0.7), Inches(5.58), Inches(12.0), Inches(0.6),
                 font_size=Pt(12), color=C_GREEN)


# ─────────────────────────────────────────────────────
# スライド5: 歴代ベスト投稿
# ─────────────────────────────────────────────────────
def slide_top_posts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, C_BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=C_BG_DARK)
    add_text_box(slide, "歴代ベスト投稿", Inches(0.5), Inches(0.25), Inches(8), Inches(0.6),
                 font_size=Pt(24), bold=True, color=C_WHITE)
    add_text_box(slide, "report_20260404.txt より", Inches(10), Inches(0.35), Inches(3), Inches(0.4),
                 font_size=Pt(12), color=C_GOLD, align=PP_ALIGN.RIGHT)

    rows = [
        ("🥇", "いいね 451件", "2023-02-01", "自己語り系",
         "「今日は誰も興味のない、自分語りです（笑）」\n美容学校時代からのキャリアを語る長文投稿",
         "リーチ720人・保存14件も同時達成", C_GREEN),
        ("🥈", "いいね 378件", "2022-08-05", "プライベート系",
         "「今さら…ほんとに今さらですが（笑）」\n2年越しで2人だけの結婚式を挙げたエピソード",
         "リーチ744人・保存5件", C_BLUE),
        ("👀", "リーチ 937人", "2024-06-18", "プライベート系",
         "「久々のプライベート😊」\n近況報告＋施術スタイルのPR",
         "いいね39件", C_GOLD),
        ("🔖", "保存 4件",    "2022-04-19", "美容情報系",
         "ヒト幹細胞培養液スキンケアの詳細解説\n「その中でも私は〜」と具体的な体験を交えた情報提供",
         "リーチ40人（フォロワー限定だが保存率高）", C_GOLD_DARK),
    ]

    for i, (medal, metric, date, category, cap, note, color) in enumerate(rows):
        ry = Inches(1.3) + i * Inches(1.45)
        add_rect(slide, Inches(0.4), ry, Inches(12.5), Inches(1.3),
                 fill_color=C_WHITE, line_color=C_LGRAY, line_width=Pt(0.5))
        add_rect(slide, Inches(0.4), ry, Inches(0.1), Inches(1.3), fill_color=color)

        # メダル・指標
        add_text_box(slide, medal + "  " + metric,
                     Inches(0.65), ry + Inches(0.1), Inches(2.5), Inches(0.45),
                     font_size=Pt(15), bold=True, color=color)
        # 日付・カテゴリ
        add_text_box(slide, f"{date}  |  {category}",
                     Inches(0.65), ry + Inches(0.5), Inches(2.8), Inches(0.35),
                     font_size=Pt(11), color=C_GRAY)
        # キャプション概要
        add_text_box(slide, cap, Inches(3.5), ry + Inches(0.05), Inches(6.2), Inches(0.75),
                     font_size=Pt(12), color=C_BLACK, wrap=True)
        # 備考
        add_text_box(slide, "📎 " + note,
                     Inches(9.7), ry + Inches(0.3), Inches(3.2), Inches(0.45),
                     font_size=Pt(11), color=C_GRAY, wrap=True)


# ─────────────────────────────────────────────────────
# スライド6: 直近10件トレンド
# ─────────────────────────────────────────────────────
def slide_recent(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, C_BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=C_BG_DARK)
    add_text_box(slide, "直近10件のトレンド（2026年3月〜4月）", Inches(0.5), Inches(0.25), Inches(10), Inches(0.6),
                 font_size=Pt(24), bold=True, color=C_WHITE)
    add_text_box(slide, "get_recent_insights.py 実行結果", Inches(10), Inches(0.35), Inches(3), Inches(0.4),
                 font_size=Pt(12), color=C_GOLD, align=PP_ALIGN.RIGHT)

    # テーブルヘッダー
    cols = ["日付", "タイプ", "テーマ", "いいね", "保存", "リーチ"]
    col_ws = [Inches(1.4), Inches(1.5), Inches(3.5), Inches(1.2), Inches(1.2), Inches(1.2)]
    col_xs = [Inches(0.5)]
    for w in col_ws[:-1]:
        col_xs.append(col_xs[-1] + w)

    header_y = Inches(1.3)
    add_rect(slide, Inches(0.5), header_y, Inches(10), Inches(0.38), fill_color=C_BG_DARK)
    for j, (label, cx) in enumerate(zip(cols, col_xs)):
        add_text_box(slide, label, cx + Inches(0.05), header_y + Inches(0.05),
                     col_ws[j] - Inches(0.1), Inches(0.3),
                     font_size=Pt(11), bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    # recent_insights.json を読む
    recent_data = []
    if os.path.exists("recent_insights.json"):
        with open("recent_insights.json", encoding="utf-8") as f:
            d = json.load(f)
        recent_data = d.get("posts", [])

    # テーブル行
    rows_data = [
        ("04/13", "カルーセル", "未分類（ブライダル推定）", 7, 0, 0),
        ("04/12", "単体画像",   "未分類",                  5, 0, 0),
        ("04/06", "カルーセル", "ライフスタイル・自己語り", 5, 0, 0),
        ("04/05", "カルーセル", "ライフスタイル・自己語り", 9, 0, 0),
        ("04/04", "単体画像",   "メニュー・サービス",       4, 0, 0),
        ("04/03", "カルーセル", "ライフスタイル・自己語り", 6, 0, 0),
        ("03/31", "単体画像",   "メニュー・サービス",      11, 0, 0),
        ("03/26", "単体画像",   "ブライダル",               6, 0, 0),
        ("03/17", "カルーセル", "ライフスタイル・自己語り", 8, 0, 0),
        ("03/13", "カルーセル", "ブライダル",              11, 0, 0),
    ]

    # recent_insights.json があれば上書き
    if recent_data:
        rows_data = []
        for p in recent_data[:10]:
            rows_data.append((
                p.get("date", "")[-5:],
                p.get("media_type", ""),
                p.get("theme", ""),
                p.get("like_count", 0),
                p.get("saved", 0),
                p.get("reach", 0),
            ))

    max_like = max(r[3] for r in rows_data) or 1

    for i, (date, mtype, theme, likes, saved, reach) in enumerate(rows_data):
        ry = header_y + Inches(0.38) + i * Inches(0.49)
        bg = C_WHITE if i % 2 == 0 else C_BG_LIGHT
        add_rect(slide, Inches(0.5), ry, Inches(10), Inches(0.46), fill_color=bg)

        like_color = C_GREEN if likes >= 10 else C_BLACK
        values = [date, mtype, theme, str(likes), str(saved) if saved else "–", str(reach) if reach else "–"]
        for j, (val, cx) in enumerate(zip(values, col_xs)):
            vc = like_color if j == 3 else C_GRAY if j in [1, 4, 5] else C_BLACK
            vb = j == 3 and likes >= 10
            add_text_box(slide, val, cx + Inches(0.05), ry + Inches(0.07),
                         col_ws[j] - Inches(0.1), Inches(0.35),
                         font_size=Pt(11), bold=vb, color=vc, align=PP_ALIGN.CENTER)

    # サマリーボックス
    sx = Inches(10.8)
    add_rect(slide, sx, Inches(1.3), Inches(2.3), Inches(5.3),
             fill_color=C_WHITE, line_color=C_GOLD, line_width=Pt(1.5))
    add_text_box(slide, "直近サマリー", sx + Inches(0.15), Inches(1.45), Inches(2.0), Inches(0.35),
                 font_size=Pt(13), bold=True, color=C_GOLD_DARK, align=PP_ALIGN.CENTER)
    add_divider(slide, sx + Inches(0.15), Inches(1.85), Inches(2.0), color=C_GOLD, thickness=Pt(1))

    summary_items = [
        ("平均いいね", "7.2件"),
        ("最多テーマ", "自己語り\n4/10件"),
        ("ブライダル", "2/10件"),
        ("メニュー", "2/10件"),
    ]
    for k, (label, val) in enumerate(summary_items):
        ky = Inches(2.0) + k * Inches(1.05)
        add_text_box(slide, label, sx + Inches(0.1), ky, Inches(2.1), Inches(0.3),
                     font_size=Pt(11), color=C_GRAY, align=PP_ALIGN.CENTER)
        add_text_box(slide, val, sx + Inches(0.1), ky + Inches(0.28), Inches(2.1), Inches(0.5),
                     font_size=Pt(18), bold=True, color=C_BLACK, align=PP_ALIGN.CENTER, wrap=True)


# ─────────────────────────────────────────────────────
# スライド7: 戦略的提言
# ─────────────────────────────────────────────────────
def slide_strategy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, C_BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=C_BG_DARK)
    add_text_box(slide, "戦略的提言 — 次の投稿に活かすポイント", Inches(0.5), Inches(0.25), Inches(12), Inches(0.6),
                 font_size=Pt(24), bold=True, color=C_WHITE)

    actions = [
        (C_GREEN,    "01", "自己語り投稿を月2〜3本必ず入れる",
         "いいね平均 29.9件（メニュー告知の 2.5倍）。\nMIKIのキャリア・体験・価値観を語る長文が最も反応が高い。"),
        (C_RED,      "02", "メニュー告知（現在54%）の比率を 40% 以下に下げる",
         "最多投稿なのに最低エンゲージメント（平均いいね 11.8件）。\n目安: 自己語り 30% / ブライダル 20% / メニュー 40% / 他 10%"),
        (C_BLUE,     "03", "プライベート投稿を月1本入れてリーチを稼ぐ",
         "リーチ最高 937人 はシンプルな近況報告投稿。\nフォロワー外へのリーチ拡大に有効。"),
        (C_GOLD_DARK,"04", "リールを月1〜2本に増やす",
         "現在67本（17%）。平均いいねはカルーセル以上。\nショートリール（30〜60秒）で新規フォロワー獲得を狙う。"),
        (C_GOLD,     "05", "キャプション冒頭にキーワードを入れる（SEO対策）",
         "1〜2行目で「30代 肌くすみ」等のテーマキーワードを自然に記述。\n3行目から「MIKIです。」で本文スタート。"),
    ]

    col_w = Inches(6.3)
    for i, (color, num, title, body) in enumerate(actions):
        col = i % 2
        row = i // 2
        ax = Inches(0.4) + col * (col_w + Inches(0.3))
        ay = Inches(1.3) + row * Inches(1.9)
        ah = Inches(1.75)

        add_rect(slide, ax, ay, col_w, ah, fill_color=C_WHITE,
                 line_color=C_LGRAY, line_width=Pt(0.5))
        add_rect(slide, ax, ay, Inches(0.55), ah, fill_color=color)

        # 番号
        add_text_box(slide, num, ax + Inches(0.0), ay + Inches(0.55), Inches(0.55), Inches(0.55),
                     font_size=Pt(13), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # タイトル
        add_text_box(slide, title, ax + Inches(0.65), ay + Inches(0.1), col_w - Inches(0.75), Inches(0.55),
                     font_size=Pt(13), bold=True, color=color, wrap=True)
        # 本文
        add_text_box(slide, body, ax + Inches(0.65), ay + Inches(0.65), col_w - Inches(0.75), Inches(0.95),
                     font_size=Pt(11), color=C_BLACK, wrap=True)


# ─────────────────────────────────────────────────────
# スライド8: チェックリスト
# ─────────────────────────────────────────────────────
def slide_checklist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, C_BG_DARK)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=C_GOLD)
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_color=C_GOLD)

    add_text_box(slide, "次回投稿 チェックリスト", Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
                 font_size=Pt(28), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_divider(slide, Inches(3), Inches(1.1), Inches(7.3), color=C_GOLD)

    checks = [
        "キャプション 1〜2行目に検索キーワードが自然に入っているか",
        "冒頭が「⭐️閲覧ありがとう」などのテンプレ告知になっていないか",
        "MIKIの体験・感情・視点が伝わる書き出しになっているか",
        "カルーセルのスライドに「保存したくなる情報」があるか",
        "スライド末尾 2枚（slide8 / slide7）は自動追加されているか",
        "絵文字の数は 3〜5個の範囲に収まっているか",
        "CTA（DM受付・初回 20%OFF）が末尾に明記されているか",
        "今月の自己語り投稿は済んでいるか（月 2〜3本ペース）",
        "スケジュール前に get_recent_insights.py を実行したか",
    ]

    col_split = 5  # 左列に5件
    for i, text in enumerate(checks):
        col = 0 if i < col_split else 1
        row = i if i < col_split else i - col_split
        cx = Inches(0.8) + col * Inches(6.4)
        cy = Inches(1.4) + row * Inches(1.0)

        # チェックボックス背景
        add_rect(slide, cx, cy, Inches(5.8), Inches(0.8),
                 fill_color=RGBColor(0x28, 0x28, 0x28),
                 line_color=C_GOLD_DARK, line_width=Pt(0.75))
        # □ マーク
        add_text_box(slide, "□", cx + Inches(0.1), cy + Inches(0.12), Inches(0.5), Inches(0.55),
                     font_size=Pt(18), color=C_GOLD, align=PP_ALIGN.CENTER)
        # テキスト
        add_text_box(slide, text, cx + Inches(0.6), cy + Inches(0.1), Inches(5.1), Inches(0.6),
                     font_size=Pt(12), color=C_WHITE, wrap=True)


# ─────────────────────────────────────────────────────
# Instagram API: オーディエンス・フォロワー取得
# ─────────────────────────────────────────────────────
def fetch_audience_insights() -> dict:
    """エリア（city）・年齢層（gender_age）データを取得。失敗時は {}。"""
    if not TOKEN or not ACCOUNT_ID:
        return {}
    try:
        url = f"{API_BASE}/{ACCOUNT_ID}/insights"
        params = {
            "metric": "audience_city,audience_gender_age",
            "period": "lifetime",
            "access_token": TOKEN,
        }
        res = requests.get(url, params=params, timeout=15)
        data = res.json()
        if "error" in data:
            print(f"  オーディエンスAPI エラー: {data['error'].get('message')}")
            return {}
        result = {}
        for item in data.get("data", []):
            name = item["name"]
            vals = (item.get("values") or [{}])[0].get("value", {})
            result[name] = vals
        return result
    except Exception as e:
        print(f"  オーディエンス取得エラー: {e}")
        return {}


def fetch_follower_trend(days: int = 30) -> list:
    """直近 days 日の新規フォロワー推移 → [{"date": "YYYY-MM-DD", "value": n}, ...]"""
    if not TOKEN or not ACCOUNT_ID:
        return []
    try:
        now = datetime.now(timezone.utc)
        since = int((now - timedelta(days=days)).timestamp())
        until = int(now.timestamp())
        url = f"{API_BASE}/{ACCOUNT_ID}/insights"
        params = {
            "metric": "follower_count",
            "period": "day",
            "since": since,
            "until": until,
            "access_token": TOKEN,
        }
        res = requests.get(url, params=params, timeout=15)
        data = res.json()
        if "error" in data:
            print(f"  フォロワーAPI エラー: {data['error'].get('message')}")
            return []
        for item in data.get("data", []):
            if item.get("name") == "follower_count":
                return [
                    {"date": v["end_time"][:10], "value": v.get("value", 0)}
                    for v in item.get("values", [])
                ]
    except Exception as e:
        print(f"  フォロワー推移取得エラー: {e}")
    return []


# ─────────────────────────────────────────────────────
# スライドA: 見てる人のエリア（都市別上位10）
# ─────────────────────────────────────────────────────
def slide_audience_location(prs, audience_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, C_BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=C_BG_DARK)
    add_text_box(slide, "フォロワーのエリア（都市別）", Inches(0.5), Inches(0.25), Inches(9), Inches(0.6),
                 font_size=Pt(24), bold=True, color=C_WHITE)
    add_text_box(slide, "audience_city / lifetime", Inches(10), Inches(0.35), Inches(3), Inches(0.4),
                 font_size=Pt(12), color=C_GOLD, align=PP_ALIGN.RIGHT)

    city_raw = audience_data.get("audience_city", {})

    if city_raw:
        sorted_cities = sorted(city_raw.items(), key=lambda x: x[1], reverse=True)[:10]
        total_shown = sum(v for _, v in sorted_cities) or 1
        max_v = sorted_cities[0][1] if sorted_cities else 1

        label_w  = Inches(2.8)
        bar_x    = Inches(3.7)
        bar_max_w = Inches(7.5)
        bar_h    = Inches(0.44)
        gap      = Inches(0.12)
        start_y  = Inches(1.3)

        colors_bar = [C_GOLD, C_GOLD, C_BLUE, C_BLUE, C_BLUE,
                      C_LGRAY, C_LGRAY, C_LGRAY, C_LGRAY, C_LGRAY]

        for i, (city, count) in enumerate(sorted_cities):
            ry = start_y + i * (bar_h + gap)
            city_name = city.replace(",JP", "").replace(",", " ").strip()

            add_text_box(slide, city_name, Inches(0.5), ry + Inches(0.04),
                         label_w, bar_h - Inches(0.08),
                         font_size=Pt(11), bold=(i < 3), color=C_BLACK, align=PP_ALIGN.RIGHT)

            fill_c = colors_bar[i]
            add_bar(slide, bar_x + Inches(0.05), ry + Inches(0.06),
                    bar_max_w - Inches(0.1), bar_h - Inches(0.12),
                    count, max_v, fill_color=fill_c, bg_color=C_LGRAY)

            pct = round(count / total_shown * 100, 1)
            label_c = C_WHITE if (count / max_v) > 0.35 else C_BLACK
            add_text_box(slide, f"  {count}人（{pct}%）",
                         bar_x + Inches(0.1), ry + Inches(0.07),
                         bar_max_w - Inches(0.2), bar_h - Inches(0.14),
                         font_size=Pt(10), bold=(i < 3), color=label_c, align=PP_ALIGN.LEFT)
    else:
        add_text_box(slide,
                     "⚠️  エリアデータを取得できませんでした\n"
                     "instagram_manage_insights 権限が必要 / アクセストークン要確認",
                     Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.2),
                     font_size=Pt(16), color=C_GRAY, align=PP_ALIGN.CENTER)

    add_divider(slide, Inches(0.5), Inches(6.9), Inches(12.3), color=C_GOLD)
    add_text_box(slide, "💡  上位エリアに合わせた投稿時間・コンテンツの地域特性を考慮する",
                 Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
                 font_size=Pt(11), color=C_GOLD_DARK)


# ─────────────────────────────────────────────────────
# スライドB: 年齢層・性別
# ─────────────────────────────────────────────────────
def slide_audience_age(prs, audience_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, C_BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=C_BG_DARK)
    add_text_box(slide, "フォロワーの年齢層・性別", Inches(0.5), Inches(0.25), Inches(9), Inches(0.6),
                 font_size=Pt(24), bold=True, color=C_WHITE)
    add_text_box(slide, "audience_gender_age / lifetime", Inches(9.0), Inches(0.35), Inches(4), Inches(0.4),
                 font_size=Pt(12), color=C_GOLD, align=PP_ALIGN.RIGHT)

    ga_raw = audience_data.get("audience_gender_age", {})

    if ga_raw:
        # キーは "F.18-24", "M.25-34" など
        age_groups = ["13-17", "18-24", "25-34", "35-44", "45-54", "55-64", "65+"]
        f_vals = {}
        m_vals = {}
        u_vals = {}
        for key, val in ga_raw.items():
            parts = key.split(".")
            if len(parts) == 2:
                gender, age = parts
                if gender == "F":
                    f_vals[age] = f_vals.get(age, 0) + val
                elif gender == "M":
                    m_vals[age] = m_vals.get(age, 0) + val
                else:
                    u_vals[age] = u_vals.get(age, 0) + val

        total = sum(f_vals.values()) + sum(m_vals.values()) + sum(u_vals.values()) or 1
        all_vals = [f_vals.get(a, 0) + m_vals.get(a, 0) + u_vals.get(a, 0) for a in age_groups]
        max_v = max(all_vals) or 1

        # 性別サマリー
        f_total = sum(f_vals.values())
        m_total = sum(m_vals.values())
        f_pct = round(f_total / total * 100, 1)
        m_pct = round(m_total / total * 100, 1)
        add_text_box(slide, f"女性 {f_pct}%  /  男性 {m_pct}%  /  不明 {round(100-f_pct-m_pct,1)}%",
                     Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.35),
                     font_size=Pt(13), bold=True, color=C_GOLD_DARK, align=PP_ALIGN.CENTER)

        # 年齢別バーチャート（女性／男性 2色積み上げ）
        label_w  = Inches(1.5)
        bar_x    = Inches(2.1)
        bar_max_w = Inches(9.0)
        bar_h    = Inches(0.52)
        gap      = Inches(0.08)
        start_y  = Inches(1.6)

        for i, age in enumerate(age_groups):
            ry = start_y + i * (bar_h + gap)
            f_v = f_vals.get(age, 0)
            m_v = m_vals.get(age, 0)
            tot_v = f_v + m_v + u_vals.get(age, 0)

            add_text_box(slide, age, Inches(0.5), ry + Inches(0.06),
                         label_w, bar_h - Inches(0.12),
                         font_size=Pt(12), bold=True, color=C_BLACK, align=PP_ALIGN.RIGHT)

            # 女性（ゴールド）
            f_w = int(bar_max_w * f_v / max_v)
            if f_w > 0:
                add_rect(slide, bar_x, ry + Inches(0.07), f_w, bar_h - Inches(0.14), fill_color=C_GOLD)
            # 男性（青）
            m_w = int(bar_max_w * m_v / max_v)
            if m_w > 0:
                add_rect(slide, bar_x + f_w, ry + Inches(0.07), m_w, bar_h - Inches(0.14), fill_color=C_BLUE)
            # 背景（残り）
            rest = bar_max_w - f_w - m_w
            if rest > 0:
                add_rect(slide, bar_x + f_w + m_w, ry + Inches(0.07), rest, bar_h - Inches(0.14), fill_color=C_LGRAY)

            pct_v = round(tot_v / total * 100, 1)
            add_text_box(slide, f"{tot_v}人（{pct_v}%）",
                         bar_x + bar_max_w + Inches(0.1), ry + Inches(0.07),
                         Inches(1.8), bar_h - Inches(0.14),
                         font_size=Pt(11), bold=(pct_v >= 20), color=C_BLACK)

        # 凡例
        add_rect(slide, Inches(2.1), Inches(6.55), Inches(0.3), Inches(0.22), fill_color=C_GOLD)
        add_text_box(slide, "女性", Inches(2.5), Inches(6.55), Inches(1.0), Inches(0.25), font_size=Pt(11), color=C_GOLD_DARK)
        add_rect(slide, Inches(3.7), Inches(6.55), Inches(0.3), Inches(0.22), fill_color=C_BLUE)
        add_text_box(slide, "男性", Inches(4.1), Inches(6.55), Inches(1.0), Inches(0.25), font_size=Pt(11), color=C_BLUE)
    else:
        add_text_box(slide,
                     "⚠️  年齢層データを取得できませんでした\n"
                     "instagram_manage_insights 権限が必要 / アクセストークン要確認",
                     Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.2),
                     font_size=Pt(16), color=C_GRAY, align=PP_ALIGN.CENTER)

    add_divider(slide, Inches(0.5), Inches(6.9), Inches(12.3), color=C_GOLD)
    add_text_box(slide, "💡  メインターゲット層に合わせたコンテンツ・言葉づかいを最適化する",
                 Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
                 font_size=Pt(11), color=C_GOLD_DARK)


# ─────────────────────────────────────────────────────
# スライドC: 新規フォロワーの推移（直近30日）
# ─────────────────────────────────────────────────────
def slide_follower_trend(prs, trend_data: list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, C_BG_LIGHT)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=C_BG_DARK)
    add_text_box(slide, "新規フォロワーの推移（直近30日）", Inches(0.5), Inches(0.25), Inches(9), Inches(0.6),
                 font_size=Pt(24), bold=True, color=C_WHITE)
    add_text_box(slide, "follower_count / day", Inches(10.5), Inches(0.35), Inches(2.5), Inches(0.4),
                 font_size=Pt(12), color=C_GOLD, align=PP_ALIGN.RIGHT)

    if trend_data:
        values = [d["value"] for d in trend_data]
        dates  = [d["date"] for d in trend_data]
        total_new = sum(values)
        avg_new   = round(total_new / len(values), 1) if values else 0
        max_v = max(values) if values else 1

        # KPI サマリー行
        kpi_items = [
            ("合計新規フォロワー", f"{total_new}人", C_GOLD),
            ("日平均",             f"{avg_new}人/日",  C_BLUE),
            ("最大（1日）",        f"{max(values)}人", C_GREEN),
        ]
        kw = Inches(3.8)
        kgap = Inches(0.35)
        ky = Inches(1.2)
        for i, (lbl, val, clr) in enumerate(kpi_items):
            kx = Inches(0.5) + i * (kw + kgap)
            add_rect(slide, kx, ky, kw, Inches(0.9), fill_color=C_WHITE,
                     line_color=clr, line_width=Pt(1.5))
            add_text_box(slide, lbl, kx + Inches(0.1), ky + Inches(0.05),
                         kw - Inches(0.2), Inches(0.28),
                         font_size=Pt(11), color=C_GRAY, align=PP_ALIGN.CENTER)
            add_text_box(slide, val, kx + Inches(0.1), ky + Inches(0.3),
                         kw - Inches(0.2), Inches(0.5),
                         font_size=Pt(24), bold=True, color=clr, align=PP_ALIGN.CENTER)

        # 折れ線グラフ（バーで代替）
        chart_x = Inches(0.5)
        chart_y = Inches(2.3)
        chart_w = Inches(12.3)
        chart_h = Inches(4.0)
        add_rect(slide, chart_x, chart_y, chart_w, chart_h,
                 fill_color=C_WHITE, line_color=C_LGRAY, line_width=Pt(0.5))

        n = len(trend_data)
        if n > 0:
            bar_w = (chart_w - Inches(0.4)) / n
            pad_x = Inches(0.2)
            pad_y = Inches(0.2)
            inner_h = chart_h - pad_y * 2

            for i, (val, date) in enumerate(zip(values, dates)):
                bx = chart_x + pad_x + i * bar_w
                ratio = val / max_v if max_v > 0 else 0
                bh = inner_h * ratio
                by = chart_y + pad_y + inner_h - bh

                bar_clr = C_GOLD if val == max(values) else C_BLUE
                add_rect(slide, int(bx), int(by), int(bar_w * 0.7), int(bh), fill_color=bar_clr)

                # 日付ラベル（7日おき）
                if i % 7 == 0 or i == n - 1:
                    label_date = date[5:]  # MM-DD
                    add_text_box(slide, label_date,
                                 int(bx) - Inches(0.15), chart_y + chart_h - Inches(0.22),
                                 Inches(0.6), Inches(0.2),
                                 font_size=Pt(8), color=C_GRAY, align=PP_ALIGN.CENTER)

                # 値ラベル（高い日のみ）
                if val >= avg_new * 1.5 and val > 0:
                    add_text_box(slide, str(val),
                                 int(bx), int(by) - Inches(0.25),
                                 int(bar_w * 0.7), Inches(0.22),
                                 font_size=Pt(9), bold=True, color=bar_clr, align=PP_ALIGN.CENTER)
    else:
        add_text_box(slide,
                     "⚠️  フォロワー推移データを取得できませんでした\n"
                     "instagram_manage_insights 権限が必要 / アクセストークン要確認",
                     Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.2),
                     font_size=Pt(16), color=C_GRAY, align=PP_ALIGN.CENTER)

    add_divider(slide, Inches(0.5), Inches(6.9), Inches(12.3), color=C_GOLD)
    add_text_box(slide, "💡  投稿日とフォロワー増加の相関を確認し、反応の良い曜日・時間帯を特定する",
                 Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
                 font_size=Pt(11), color=C_GOLD_DARK)


# ─────────────────────────────────────────────────────
# Google Drive アップロード
# ─────────────────────────────────────────────────────
def upload_pptx_to_drive(local_path: str, filename: str) -> str:
    """PPTXをGoogle Driveの指定フォルダにアップロードしてURLを返す（OAuth2ユーザー認証）"""
    try:
        from google.oauth2.credentials import Credentials as UserCredentials
        from google.auth.transport.requests import Request
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaFileUpload

        TOKEN_PATH = "token_drive.json"
        with open(TOKEN_PATH, "r") as f:
            info = json.load(f)
        creds = UserCredentials(
            token=info["token"],
            refresh_token=info["refresh_token"],
            token_uri=info["token_uri"],
            client_id=info["client_id"],
            client_secret=info["client_secret"],
            scopes=info["scopes"],
        )
        if creds.expired and creds.refresh_token:
            creds.refresh(Request())
            with open(TOKEN_PATH, "w") as f:
                f.write(creds.to_json())
        service = build("drive", "v3", credentials=creds)

        # drive_folders.json から lifestyle フォルダIDを取得（サービスアカウントはMy Driveに書けないため共有フォルダを使用）
        folder_id = os.getenv("DRIVE_FOLDER_ID")
        if not folder_id:
            try:
                with open("drive_folders.json", "r", encoding="utf-8") as f:
                    folder_id = json.load(f).get("lifestyle", "")
            except Exception:
                folder_id = ""

        file_metadata = {"name": filename}
        if folder_id:
            file_metadata["parents"] = [folder_id]

        mimetype = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        media = MediaFileUpload(local_path, mimetype=mimetype)

        # 同名ファイルがあれば上書き
        q = f"name='{filename}' and trashed=false"
        if folder_id:
            q += f" and '{folder_id}' in parents"
        existing = service.files().list(q=q, fields="files(id)").execute().get("files", [])

        if existing:
            file_id = existing[0]["id"]
            service.files().update(fileId=file_id, media_body=media).execute()
        else:
            result = service.files().create(
                body=file_metadata, media_body=media, fields="id"
            ).execute()
            file_id = result.get("id", "")

        # 誰でも閲覧可能に（オプション）
        service.permissions().create(
            fileId=file_id,
            body={"type": "anyone", "role": "reader"},
        ).execute()

        url = f"https://drive.google.com/file/d/{file_id}/view"
        return url
    except Exception as e:
        print(f"  Drive アップロードエラー: {e}")
        return ""


# ─────────────────────────────────────────────────────
# メイン
# ─────────────────────────────────────────────────────
def main():
    print("PowerPoint レポート生成中...")

    # Instagram API からオーディエンス・フォロワーデータを取得
    print("  Instagram API からオーディエンスデータを取得中...")
    audience_data = fetch_audience_insights()
    print(f"  取得済みメトリクス: {list(audience_data.keys()) or '(なし)'}")

    print("  フォロワー推移データを取得中...")
    trend_data = fetch_follower_trend(days=30)
    print(f"  フォロワー推移: {len(trend_data)}日分")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("  スライド 1: タイトル")
    slide_title(prs)
    print("  スライド 2: アカウントサマリー")
    slide_summary(prs)
    print("  スライド 3: テーマ別パフォーマンス")
    slide_theme_perf(prs)
    print("  スライド 4: コンテンツタイプ別")
    slide_type_perf(prs)
    print("  スライド 5: 歴代ベスト投稿")
    slide_top_posts(prs)
    print("  スライド 6: 直近10件トレンド")
    slide_recent(prs)
    print("  スライド 7: フォロワーのエリア")
    slide_audience_location(prs, audience_data)
    print("  スライド 8: フォロワーの年齢層・性別")
    slide_audience_age(prs, audience_data)
    print("  スライド 9: 新規フォロワーの推移")
    slide_follower_trend(prs, trend_data)
    print("  スライド 10: 戦略的提言")
    slide_strategy(prs)
    print("  スライド 11: チェックリスト")
    slide_checklist(prs)

    prs.save(OUTPUT)
    print(f"\n✅ 保存完了: {OUTPUT}")
    print(f"   → Finder で開くか: open '{OUTPUT}'")

    # Google Drive にアップロード
    print("\nGoogle Drive にアップロード中...")
    drive_url = upload_pptx_to_drive(OUTPUT, "instagram_insight_latest.pptx")
    if drive_url:
        print(f"✅ Drive アップロード完了: {drive_url}")
    else:
        print("  Drive アップロードはスキップ（credentials.json が必要）")


if __name__ == "__main__":
    main()
